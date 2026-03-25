"""
知识库解析：PDF/PPT/Excel/Word
仅用于 AI 回复参考，不可作为附件外发
"""
import os
from pathlib import Path
from typing import Optional

# PDF
from pypdf import PdfReader

# PPT
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Excel
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

# Word
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


def parse_pdf(file_path: str) -> str:
    """解析 PDF 文本"""
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            texts.append(t)
    return "\n\n".join(texts)


def parse_pptx(file_path: str) -> str:
    """解析 PPT 文本"""
    if not Presentation:
        raise ImportError("请安装 python-pptx")
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n\n".join(texts)


def parse_xlsx(file_path: str) -> str:
    """解析 Excel 文本"""
    if not load_workbook:
        raise ImportError("请安装 openpyxl")
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(c) for c in row if c)
            if row_text:
                texts.append(row_text)
    return "\n\n".join(texts)


def parse_docx(file_path: str) -> str:
    """解析 Word 文本"""
    if not DocxDocument:
        raise ImportError("请安装 python-docx")
    doc = DocxDocument(file_path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text)


def parse_file(file_path: str, file_type: Optional[str] = None) -> str:
    """
    根据扩展名解析文件
    支持: pdf, pptx, xlsx, docx
    """
    ext = (file_type or Path(file_path).suffix).lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    if ext in (".pptx", ".ppt"):
        return parse_pptx(file_path)
    if ext in (".xlsx", ".xls"):
        return parse_xlsx(file_path)
    if ext in (".docx", ".doc"):
        return parse_docx(file_path)
    raise ValueError(f"不支持的文件格式: {ext}")
